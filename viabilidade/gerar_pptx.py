import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Definições de Cores (Tema Dark/Tech Minimalista)
BG_COLOR = RGBColor(15, 23, 42)      # Slate 900 (#0f172a)
TITLE_COLOR = RGBColor(56, 189, 248) # Sky 400 (#38bdf8)
TEXT_COLOR = RGBColor(203, 213, 225) # Slate 300 (#cbd5e1)
ACCENT_COLOR = RGBColor(248, 113, 113)# Red 400 (#f87171)

FONT_NAME = 'Plus Jakarta Sans'

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title(slide, text, top=Inches(0.5)):
    title_box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.33), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    return title_box

def add_text(slide, text, left, top, width, height, font_size=20, bold=False, color=TEXT_COLOR):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return txBox

def add_image_centered(slide, image_path, top_offset=Inches(2)):
    if os.path.exists(image_path):
        # Center horizontally, place below title
        # Slide width is 13.33 inches in 16:9
        # Assuming we want the image to be around 8 inches wide
        width = Inches(8)
        left = (Inches(13.33) - width) / 2
        slide.shapes.add_picture(image_path, left, top_offset, width=width)

def criar_apresentacao():
    prs = Presentation()
    # 16:9 Widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    # --- SLIDE 1: Capa ---
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide)
    
    add_text(slide, "Marketplace B2B", Inches(1), Inches(2.5), Inches(11.33), Inches(1), font_size=54, bold=True, color=TITLE_COLOR)
    add_text(slide, "Construção Civil: Eficiência & Transparência", Inches(1), Inches(3.5), Inches(11.33), Inches(1), font_size=32, bold=True, color=TEXT_COLOR)
    add_text(slide, "Análise Institucional e Viabilidade de Risco (Seed Stage)", Inches(1), Inches(4.5), Inches(11.33), Inches(1), font_size=20, color=TEXT_COLOR)
    
    # --- SLIDE 2: Problema ---
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide)
    add_title(slide, "O Problema")
    add_text(slide, "A cadeia de suprimentos da construção é um ambiente de alta fricção e extrema ineficiência.", Inches(1), Inches(2), Inches(11), Inches(1), font_size=24)
    add_text(slide, "• Processos de cotação 100% analógicos (WhatsApp/Telefone).\n• Forte assimetria de informações e preços.\n• Custo elevado na descoberta de novos fornecedores.", Inches(1), Inches(3), Inches(11), Inches(3), font_size=22)

    # --- SLIDE 3: Solução ---
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide)
    add_title(slide, "O Modelo (Premissas Reais)")
    add_text(slide, "O modelo não subestima o risco operacional da tese.", Inches(1), Inches(2), Inches(11), Inches(1), font_size=24)
    add_text(slide, "• Monetização: Take rate de 2,5% sobre o GMV.\n• Risco de Bypass: Assumimos desintermediação de 50%. Capturamos apenas metade do potencial de materiais da obra.\n• Fricção de Vendas: Conversão espremida em agressivos 5% (Cold outreach B2B).", Inches(1), Inches(3), Inches(11), Inches(3), font_size=22)

    # --- SLIDE 4: Mercado ---
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide)
    add_title(slide, "A Escala do Jogo (TAM/SAM/SOM)")
    
    # Cards
    add_text(slide, "R$ 64 B", Inches(1), Inches(3), Inches(3), Inches(1), font_size=40, bold=True, color=TITLE_COLOR)
    add_text(slide, "TAM (Total Brasil)", Inches(1), Inches(4), Inches(3), Inches(0.5), font_size=16)
    
    add_text(slide, "R$ 12,8 B", Inches(5), Inches(3), Inches(3), Inches(1), font_size=40, bold=True, color=TITLE_COLOR)
    add_text(slide, "SAM (Mercado Formal)", Inches(5), Inches(4), Inches(3), Inches(0.5), font_size=16)
    
    add_text(slide, "1.264", Inches(9), Inches(3), Inches(3), Inches(1), font_size=40, bold=True, color=TITLE_COLOR)
    add_text(slide, "Obras (Meta SOM Ano 5)", Inches(9), Inches(4), Inches(3), Inches(0.5), font_size=16)

    # --- SLIDE 5: Unit Economics ---
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide)
    add_title(slide, "Unit Economics vs Custo de Vendas")
    add_text(slide, "Com conversão de apenas 5%, o CAC atinge até R$ 15.000. Contudo, a alta margem do Ticket defende a eficiência do capital.", Inches(0.5), Inches(1.2), Inches(12.33), Inches(1), font_size=20)
    add_image_centered(slide, "relatorio_latex/figuras/grafico_unit_economics.png", top_offset=Inches(2.0))

    # --- SLIDE 6: Solver ---
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide)
    add_title(slide, "Auditoria Orçamentária (Simplex)")
    add_text(slide, "A otimização linear atesta que a verba restrita de marketing suporta a meta, sendo desafiada apenas a partir do Ano 4.", Inches(0.5), Inches(1.2), Inches(12.33), Inches(1), font_size=20)
    add_image_centered(slide, "relatorio_latex/figuras/grafico_solver_mix.png", top_offset=Inches(2.0))

    # --- SLIDE 7: FCL ---
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide)
    add_title(slide, "O 'Vale da Morte' e Fluxo de Caixa")
    add_text(slide, "O Investimento de R$ 705 Mil cobre 6 meses de inoperância comercial e queima do OPEX inicial.", Inches(0.5), Inches(1.2), Inches(12.33), Inches(1), font_size=20)
    add_image_centered(slide, "relatorio_latex/figuras/grafico_fcl.png", top_offset=Inches(2.0))

    # --- SLIDE 8: Indicadores ---
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide)
    add_title(slide, "Indicadores Chave de Performance")
    add_text(slide, "O rigor matemático expõe a atratividade do ativo:", Inches(1), Inches(1.5), Inches(11), Inches(1), font_size=22)
    
    add_text(slide, "R$ 12,59 M", Inches(1), Inches(2.5), Inches(3), Inches(1), font_size=36, bold=True, color=TITLE_COLOR)
    add_text(slide, "VPL Base", Inches(1), Inches(3.2), Inches(3), Inches(0.5), font_size=16)
    
    add_text(slide, "319%", Inches(5), Inches(2.5), Inches(3), Inches(1), font_size=36, bold=True, color=TITLE_COLOR)
    add_text(slide, "TIR", Inches(5), Inches(3.2), Inches(3), Inches(0.5), font_size=16)
    
    add_text(slide, "9,5 meses", Inches(9), Inches(2.5), Inches(3), Inches(1), font_size=36, bold=True, color=TITLE_COLOR)
    add_text(slide, "Payback Simples", Inches(9), Inches(3.2), Inches(3), Inches(0.5), font_size=16)

    add_text(slide, "45,49%", Inches(1), Inches(4.5), Inches(3), Inches(1), font_size=36, bold=True, color=ACCENT_COLOR)
    add_text(slide, "WACC (Custo de Capital)", Inches(1), Inches(5.2), Inches(3), Inches(0.5), font_size=16)
    
    add_text(slide, "14%", Inches(5), Inches(4.5), Inches(3), Inches(1), font_size=36, bold=True, color=TITLE_COLOR)
    add_text(slide, "Breakeven Conversão", Inches(5), Inches(5.2), Inches(3), Inches(0.5), font_size=16)
    
    add_text(slide, "> 14x", Inches(9), Inches(4.5), Inches(3), Inches(1), font_size=36, bold=True, color=TITLE_COLOR)
    add_text(slide, "LTV / CAC", Inches(9), Inches(5.2), Inches(3), Inches(0.5), font_size=16)

    # --- SLIDE 9: Motores ---
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide)
    add_title(slide, "Os Motores do Resultado")
    add_text(slide, "Como o projeto sobrevive a tantas premissas punitivas?", Inches(1), Inches(1.5), Inches(11), Inches(1), font_size=22)
    add_text(slide, "• O Poder do Ticket B2B: Mesmo com fuga transacional de 50%, a intermediação de materiais para uma única obra gera um fee absoluto superior a R$ 16 mil.\n\n• OPEX Enxuto e Escalável: Sem frota logística. O custo central (folha de pagamentos) dilui com o volume (Margem EBITDA atinge 83% no Ano 5).\n\n• Defesa Comercial: Mesmo que o CAC do Outbound dispare para R$ 15 mil (conversão de 5%), o LTV contínuo de R$ 120 mil defende agressivamente o investimento inicial.", Inches(1), Inches(2.5), Inches(11), Inches(4), font_size=18)

    # --- SLIDE 10: Monte Carlo ---
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide)
    add_title(slide, "Gestão de Risco Estocástico (Monte Carlo)")
    add_text(slide, "2.000 iterações variando Conversão (2% a 8%) e Custo de Capital (Até 70% a.a).", Inches(0.5), Inches(1.2), Inches(12.33), Inches(1), font_size=20)
    add_image_centered(slide, "relatorio_latex/figuras/grafico_monte_carlo.png", top_offset=Inches(2.0))

    # --- SLIDE 11: Conclusão ---
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide)
    add_title(slide, "Veredito Institucional")
    add_text(slide, "PROJETO APROVADO E BLINDADO", Inches(1), Inches(2), Inches(11), Inches(1), font_size=32, bold=True, color=TITLE_COLOR)
    add_text(slide, "Mesmo drenando as premissas com bypass de 50%, gateway de pagamento corroendo receita e gargalo comercial severo (5% Conv.), o modelo gera Caixa Livre persistente e sustentável.", Inches(1), Inches(3.5), Inches(11), Inches(1.5), font_size=22)
    add_text(slide, "A assimetria entre as despesas limitadas do software (Asset Light) contra a profundidade financeira colossal do setor construtivo faz do projeto uma tese elegível para Captação Privada.", Inches(1), Inches(5), Inches(11), Inches(1.5), font_size=22)

    # Save
    out_path = "apresentacao_pitch.pptx"
    prs.save(out_path)
    print(f"Apresentacao salva em {out_path}")

if __name__ == '__main__':
    criar_apresentacao()
